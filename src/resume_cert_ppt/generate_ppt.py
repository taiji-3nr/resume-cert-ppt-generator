from __future__ import annotations

import argparse
import json
from collections import defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


NAVY = RGBColor(16, 34, 53)
TEAL = RGBColor(57, 165, 169)
AMBER = RGBColor(196, 154, 44)
INK = RGBColor(36, 55, 70)
MUTED = RGBColor(72, 101, 117)
WHITE = RGBColor(255, 255, 255)
PALE_BLUE = RGBColor(244, 248, 250)
PALE_GREEN = RGBColor(241, 247, 244)
PALE_AMBER = RGBColor(248, 245, 236)
PALE_VIOLET = RGBColor(245, 242, 248)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: int = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.LEFT,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, line in enumerate(text.splitlines() or [""]):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = line
        para.alignment = align
        para.font.name = "Yu Gothic"
        para.font.size = Pt(size)
        para.font.bold = bold
        para.font.color.rgb = color
    return box


def add_panel(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: RGBColor,
    line: RGBColor | None = None,
    size: int = 16,
    color: RGBColor = INK,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
    frame = shape.text_frame
    frame.clear()
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    frame.word_wrap = True
    for index, part in enumerate(text.splitlines() or [""]):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.text = part
        para.alignment = align
        para.font.name = "Yu Gothic"
        para.font.size = Pt(size)
        para.font.bold = bold
        para.font.color.rgb = color
    return shape


def add_background(slide, image_path: Path | None):
    if image_path and image_path.exists():
        slide.shapes.add_picture(str(image_path), 0, 0, width=Inches(13.333), height=Inches(7.5))
    else:
        rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        rect.fill.solid()
        rect.fill.fore_color.rgb = RGBColor(244, 248, 250)
        rect.line.fill.background()


def cert_label(cert: dict[str, Any]) -> str:
    month = cert["month"]
    date = f"{cert['year']}.{month:02d}" if isinstance(month, int) else f"{cert['year']}.{month}"
    name = cert["name"]
    if cert.get("current_name"):
        name = f"{name}\n(現: {cert['current_name']})"
    return f"{date}\n{name}"


def build_presentation(data: dict[str, Any], image_path: Path | None, output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    certs = data["certifications"]
    person = data["person"]
    theme = data["theme"]

    slide = prs.slides.add_slide(blank)
    add_background(slide, image_path)
    add_panel(slide, theme["title"], 0.67, 0.58, 7.1, 1.22, WHITE, size=34, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
    add_text(slide, f"{person['name']} | {theme['period']}", 0.85, 1.75, 6.6, 0.5, size=20, color=RGBColor(43, 93, 112), bold=True)
    add_panel(slide, theme["summary"], 0.85, 2.75, 5.8, 1.05, NAVY, line=TEAL, size=16, color=WHITE, align=PP_ALIGN.LEFT)

    slide = prs.slides.add_slide(blank)
    add_text(slide, "基礎技術を固めた時期", 0.6, 0.35, 10.0, 0.6, size=30, color=NAVY, bold=True)
    add_panel(slide, "", 0.9, 2.72, 11.35, 0.06, TEAL)
    early = [certs[0], certs[1], certs[2]]
    positions = [(0.9, 1.45), (4.0, 3.45), (7.3, 1.45)]
    for cert, (x, y) in zip(early, positions):
        label = cert_label(cert).splitlines()
        add_panel(slide, label[0], x, y, 1.8, 0.48, NAVY, size=18, color=WHITE, bold=True)
        add_panel(slide, "\n".join(label[1:]), x, y + 0.62, 2.8, 0.9, WHITE, line=RGBColor(217, 227, 234), size=13)
    add_text(slide, "ITキャリア初期から中堅期にかけて、業務システム開発の土台となる国家資格・専門資格を取得。", 0.6, 6.32, 10.8, 0.4, size=14, color=MUTED)

    slide = prs.slides.add_slide(blank)
    add_text(slide, "教育・指導領域への展開", 0.6, 0.35, 10.0, 0.6, size=30, color=NAVY, bold=True)
    grouped: dict[str, list[str]] = defaultdict(list)
    for cert in certs[3:]:
        grouped[cert["category"]].append(f"{cert['year']} {cert['name']}")
    columns = [
        ("Office実務", PALE_BLUE, 0.7, grouped["Office実務"]),
        ("IT基礎・業務理解", PALE_AMBER, 3.9, grouped["IT基礎・業務理解"]),
        ("プログラミング教育", PALE_GREEN, 7.1, grouped["プログラミング教育"]),
        ("指導・教育", PALE_VIOLET, 10.3, grouped["指導・教育"]),
    ]
    for title, fill, x, items in columns:
        add_panel(slide, title, x, 1.34, 2.55, 0.55, fill, line=RGBColor(213, 227, 234), size=17, bold=True)
        add_text(slide, "\n".join(items), x + 0.15, 2.1, 2.25, 2.7, size=12, color=INK)
    add_panel(slide, "セカンドキャリアでは、使える力を教える力へ転換。資格の幅がそのまま指導対象の広がりを示している。", 1.25, 6.02, 10.75, 0.7, NAVY, size=14, color=WHITE)

    slide = prs.slides.add_slide(blank)
    add_text(slide, "スキル領域マップ", 0.6, 0.35, 10.0, 0.6, size=30, color=NAVY, bold=True)
    add_panel(slide, "実務 × 教育", 5.1, 2.72, 3.1, 0.85, NAVY, line=TEAL, size=23, color=WHITE, bold=True)
    add_panel(slide, "IT基礎・業務理解\n第二種情報処理 / シスアド / P検", 1.0, 1.35, 3.3, 1.0, WHITE, line=TEAL, size=14)
    add_panel(slide, "専門技術\nデータベーススペシャリスト", 8.75, 1.35, 3.3, 1.0, WHITE, line=TEAL, size=14)
    add_panel(slide, "Office実務\nMOS / Excel / Access / PowerPoint", 1.0, 4.95, 3.3, 1.0, WHITE, line=AMBER, size=14)
    add_panel(slide, "指導・教育\nICT支援員 / Scratch / 認定インストラクター", 8.75, 4.95, 3.3, 1.0, WHITE, line=AMBER, size=14)
    add_text(slide, "資格群は、40年のIT実務経験とセカンドキャリアの教育実践を接続する証跡。", 4.7, 4.2, 4.1, 0.5, size=13, color=MUTED, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    add_background(slide, image_path)
    add_panel(slide, "資格取得のストーリー", 0.7, 0.65, 6.25, 5.65, WHITE, size=30, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
    add_text(
        slide,
        "1. IT基礎を国家資格で体系化\n2. DBなど専門技術で開発マネジメントを補強\n3. MOS・サーティファイでOffice実務力を可視化\n4. ICT支援・Scratch・認定講師で教育領域へ展開",
        1.0,
        1.75,
        5.55,
        3.0,
        size=16,
        color=INK,
    )
    add_panel(slide, theme["message"], 1.0, 5.35, 5.45, 0.65, NAVY, line=TEAL, size=20, color=WHITE, bold=True)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate a certification history PowerPoint.")
    parser.add_argument("--data", type=Path, default=Path("data/certifications.json"))
    parser.add_argument("--image", type=Path, default=Path("out/assets/certification-history-bg.png"))
    parser.add_argument("--output", type=Path, default=Path("out/資格取得ヒストリー_太地稔_20260103_python.pptx"))
    args = parser.parse_args()

    data = json.loads(args.data.read_text(encoding="utf-8"))
    build_presentation(data, args.image, args.output)
    print(args.output)


if __name__ == "__main__":
    main()
